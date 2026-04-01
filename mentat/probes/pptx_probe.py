from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from mentat.probes.base import (
    BaseProbe,
    ProbeResult,
    TopicInfo,
    StructureInfo,
    TocEntry,
    Chunk,
)
from mentat.probes._utils import estimate_tokens
from mentat.librarian.instruction_templates import (
    PPTX_BRIEF_INTRO,
    PPTX_INSTRUCTIONS,
)

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False


class PPTXProbe(BaseProbe):
    """Probe for PowerPoint presentations (.pptx)."""

    def can_handle(self, filename: str, content_type: str) -> bool:
        if not _PPTX_AVAILABLE:
            return False
        return filename.lower().endswith(".pptx")

    def run(self, file_path: str, **kwargs) -> ProbeResult:
        prs = Presentation(file_path)

        toc_entries: List[TocEntry] = []
        chunks: List[Chunk] = []
        total_words = 0
        total_images = 0
        has_notes = False

        for i, slide in enumerate(prs.slides):
            slide_num = i + 1

            # --- Slide title ---
            title_text = None
            if slide.shapes.title:
                title_text = slide.shapes.title.text.strip()

            # --- Body text (bullets) ---
            bullets: List[str] = []
            slide_text_parts: List[str] = []
            image_count = 0
            table_count = 0

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_text_parts.append(text)
                            # Non-title text as bullets
                            if shape != slide.shapes.title:
                                bullets.append(text)

                if hasattr(shape, "shape_type"):
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_count += 1
                    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        table_count += 1

                if shape.has_table:
                    table_count = max(table_count, 1)

            total_images += image_count

            # --- Speaker notes ---
            notes_text = None
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame and notes_frame.text.strip():
                    notes_text = notes_frame.text.strip()
                    has_notes = True

            # --- Annotation ---
            annot_parts = []
            if bullets:
                annot_parts.append(f"{len(bullets)} bullets")
            if image_count:
                annot_parts.append(f"{image_count} images")
            if table_count:
                annot_parts.append(f"{table_count} tables")
            annotation = " | ".join(annot_parts) if annot_parts else None

            # --- Preview from notes or first bullet ---
            preview = None
            if notes_text:
                preview = notes_text[:120]
            elif bullets:
                preview = bullets[0][:120]

            # --- ToC entry for slide ---
            slide_title = f"Slide {slide_num}: {title_text}" if title_text else f"Slide {slide_num}"
            toc_entries.append(
                TocEntry(level=1, title=slide_title, annotation=annotation, preview=preview)
            )

            # --- Bullet sub-entries (first 3 per slide) ---
            for bullet in bullets[:3]:
                toc_entries.append(
                    TocEntry(level=2, title=bullet[:80])
                )

            # --- Chunk for this slide ---
            chunk_parts = []
            if title_text:
                chunk_parts.append(f"# {title_text}")
            chunk_parts.extend(slide_text_parts)
            if notes_text:
                chunk_parts.append(f"\n[Notes] {notes_text}")

            slide_content = "\n".join(chunk_parts)
            total_words += len(slide_content.split())

            chunks.append(
                Chunk(
                    content=slide_content,
                    index=i,
                    section=slide_title,
                )
            )

        # --- Stats ---
        approx_tokens = estimate_tokens(
            "\n".join(c.content for c in chunks)
        )
        stats: Dict[str, Any] = {
            "slide_count": len(prs.slides),
            "total_word_count": total_words,
            "total_image_count": total_images,
            "has_notes": has_notes,
            "approx_tokens": approx_tokens,
        }

        # --- Topic ---
        first_title = None
        for e in toc_entries:
            if e.level == 1:
                first_title = e.title
                break

        topic = TopicInfo(
            title=first_title or Path(file_path).stem,
            first_paragraph=(
                f"PowerPoint presentation with {len(prs.slides)} slides, "
                f"{total_words} words, {total_images} images"
            ),
        )

        structure = StructureInfo(toc=toc_entries)

        result = ProbeResult(
            filename=Path(file_path).name,
            file_type="pptx",
            topic=topic,
            structure=structure,
            stats=stats,
            chunks=chunks,
            raw_snippet=chunks[0].content[:500] if chunks else None,
        )

        # Generate format-specific instructions
        brief_intro, instructions = self.generate_instructions(result)
        result.brief_intro = brief_intro
        result.instructions = instructions

        return result

    def generate_instructions(self, probe_result: ProbeResult) -> Tuple[str, str]:
        """Generate PPTX-specific instructions."""
        # Brief intro
        brief_intro = PPTX_BRIEF_INTRO

        # Full instructions
        instructions = PPTX_INSTRUCTIONS.format(filename=probe_result.filename)

        return brief_intro, instructions
